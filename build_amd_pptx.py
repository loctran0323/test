from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


OUTPUT_PATH = Path("/workspace/amd-presentation-draft.pptx")

BACKGROUND = RGBColor(15, 23, 42)
PANEL = RGBColor(17, 24, 39)
ACCENT = RGBColor(249, 115, 22)
ACCENT_LIGHT = RGBColor(253, 186, 116)
TEXT = RGBColor(229, 231, 235)
MUTED = RGBColor(148, 163, 184)
WHITE = RGBColor(255, 255, 255)
LINE = RGBColor(51, 65, 85)


def set_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BACKGROUND


def add_header(slide, title, kicker=None):
    if kicker:
        kicker_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(5.2), Inches(0.3))
        frame = kicker_box.text_frame
        frame.clear()
        p = frame.paragraphs[0]
        run = p.add_run()
        run.text = kicker.upper()
        run.font.name = "Arial"
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.color.rgb = ACCENT_LIGHT

    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.65), Inches(11.8), Inches(0.75))
    frame = title_box.text_frame
    frame.clear()
    frame.word_wrap = True
    p = frame.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = "Arial"
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.color.rgb = WHITE

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.42), Inches(2.6), Inches(0.08))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT
    accent.line.fill.background()


def add_footer(slide, slide_number):
    footer_box = slide.shapes.add_textbox(Inches(0.6), Inches(6.95), Inches(12.0), Inches(0.25))
    frame = footer_box.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"Group 7 | AMD presentation | Slide {slide_number}"
    run.font.name = "Arial"
    run.font.size = Pt(10)
    run.font.color.rgb = MUTED


def add_panel(slide, left, top, width, height, title=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = PANEL
    shape.line.color.rgb = LINE
    shape.line.width = Pt(1.0)
    shape.adjustments[0] = 0.08

    if title:
        title_box = slide.shapes.add_textbox(left + Inches(0.18), top + Inches(0.12), width - Inches(0.36), Inches(0.3))
        frame = title_box.text_frame
        frame.clear()
        p = frame.paragraphs[0]
        run = p.add_run()
        run.text = title
        run.font.name = "Arial"
        run.font.size = Pt(15)
        run.font.bold = True
        run.font.color.rgb = ACCENT_LIGHT

    return shape


def add_bullets(slide, left, top, width, height, bullets, font_size=18):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.NONE
    frame.clear()

    for idx, text in enumerate(bullets):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.level = 0
        p.space_after = Pt(6)
        p.bullet = True
        run = p.add_run()
        run.text = text
        run.font.name = "Arial"
        run.font.size = Pt(font_size)
        run.font.color.rgb = TEXT


def add_text_block(slide, left, top, width, height, text, font_size=17, color=TEXT, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    frame.clear()
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Arial"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_metric(slide, left, top, width, height, value, label):
    add_panel(slide, left, top, width, height)
    add_text_block(slide, left + Inches(0.16), top + Inches(0.14), width - Inches(0.32), Inches(0.45), value, font_size=24, color=WHITE, bold=True)
    add_text_block(slide, left + Inches(0.16), top + Inches(0.7), width - Inches(0.32), height - Inches(0.82), label, font_size=12, color=MUTED)


def add_source(slide, text):
    add_text_block(slide, Inches(0.85), Inches(6.45), Inches(11.2), Inches(0.28), text, font_size=10, color=MUTED)


def title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_panel(slide, Inches(0.55), Inches(0.45), Inches(12.2), Inches(6.2))

    add_text_block(
        slide,
        Inches(0.95),
        Inches(0.95),
        Inches(8.7),
        Inches(1.15),
        "Acid Mine Drainage (AMD): An Environmental Challenge in Mining Supply Chains",
        font_size=28,
        color=WHITE,
        bold=True,
    )
    add_text_block(
        slide,
        Inches(0.98),
        Inches(2.02),
        Inches(7.6),
        Inches(1.05),
        "CEE 304 presentation focused on how mining chemistry creates long-term water pollution and cleanup obligations.",
        font_size=18,
    )

    info_panel = add_panel(slide, Inches(8.85), Inches(0.95), Inches(2.95), Inches(2.55), "Class details")
    info_panel.fill.fore_color.rgb = RGBColor(32, 21, 12)
    info_panel.line.color.rgb = ACCENT
    add_bullets(
        slide,
        Inches(9.1),
        Inches(1.42),
        Inches(2.45),
        Inches(1.7),
        [
            "Group 7",
            "Professor Peters",
            "CEE 304",
            "Apr. 13, 2026",
        ],
        font_size=15,
    )

    add_panel(slide, Inches(0.95), Inches(3.25), Inches(6.0), Inches(2.35), "Presentation roadmap")
    add_bullets(
        slide,
        Inches(1.18),
        Inches(3.7),
        Inches(5.55),
        Inches(1.65),
        [
            "Challenge and supply-chain connection",
            "Science, engineering, and a simple neutralization calculation",
            "Iron Mountain Mine case study",
            "Solutions, economics, and regulation",
        ],
        font_size=16,
    )

    add_panel(slide, Inches(7.25), Inches(3.75), Inches(4.55), Inches(1.85), "Key message")
    add_text_block(
        slide,
        Inches(7.5),
        Inches(4.2),
        Inches(4.05),
        Inches(1.1),
        "AMD is not just dirty water at one mine site. It is a built-in supply-chain risk that can require decades of treatment.",
        font_size=16,
        color=WHITE,
        bold=True,
    )

    add_footer(slide, 1)


def bullet_slide(prs, number, kicker, title, bullets, source_text):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header(slide, title, kicker)
    add_panel(slide, Inches(0.6), Inches(1.7), Inches(12.1), Inches(4.95))
    add_bullets(slide, Inches(0.92), Inches(2.02), Inches(11.4), Inches(4.25), bullets, font_size=19)
    add_source(slide, source_text)
    add_footer(slide, number)


def chemistry_slide(prs, number):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header(slide, "Science and engineering principles", "Chemistry")

    add_panel(slide, Inches(0.7), Inches(1.85), Inches(6.05), Inches(4.75), "How AMD forms")
    add_bullets(
        slide,
        Inches(0.95),
        Inches(2.33),
        Inches(5.55),
        Inches(3.95),
        [
            "Pyrite reacts with oxygen and water after mining exposes sulfide-bearing rock.",
            "The reaction generates sulfate, dissolved iron, and hydrogen ions, which lower pH.",
            "Ferric iron can keep attacking more pyrite, so the process can become self-propagating.",
            "Iron- and sulfur-oxidizing microbes accelerate the chemistry.",
        ],
        font_size=16,
    )

    add_panel(slide, Inches(6.95), Inches(1.85), Inches(5.25), Inches(1.45), "Simplified reaction")
    add_text_block(
        slide,
        Inches(7.22),
        Inches(2.35),
        Inches(4.7),
        Inches(0.46),
        "2FeS2 + 7O2 + 2H2O -> 2Fe2+ + 4SO4^2- + 4H+",
        font_size=18,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    add_panel(slide, Inches(6.95), Inches(3.55), Inches(5.25), Inches(3.05), "Why low pH matters")
    add_bullets(
        slide,
        Inches(7.18),
        Inches(4.0),
        Inches(4.75),
        Inches(2.2),
        [
            "Low pH keeps metals such as Fe, Al, Cu, Zn, Pb, and Mn dissolved.",
            "When pH rises later, iron hydroxides can precipitate and coat streambeds.",
            "Those orange deposits, often called yellowboy, damage habitat and stream ecology.",
        ],
        font_size=15,
    )

    add_source(slide, "Sources: Cox et al. (2003); OSMRE (2000); U.S. Geological Survey AMD overview; Zipper et al. (2023).")
    add_footer(slide, number)


def calculation_slide(prs, number):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header(slide, "Simple calculation: why treatment burdens get large fast", "Calculation")

    add_metric(slide, Inches(0.8), Inches(1.82), Inches(2.85), Inches(1.45), "4 mol H+", "Approximate acidity produced per mole of pyrite after combining the oxidation sequence.")
    add_metric(slide, Inches(3.95), Inches(1.82), Inches(2.85), Inches(1.45), "2 mol CaCO3", "Neutralization needed for each mole of pyrite because CaCO3 consumes 2 moles of H+.")
    add_metric(slide, Inches(7.1), Inches(1.82), Inches(2.85), Inches(1.45), "1.67 kg CaCO3 / kg pyrite", "Neutralization requirement using about 120 g/mol for pyrite and 100 g/mol for CaCO3.")
    add_metric(slide, Inches(10.25), Inches(1.82), Inches(2.05), Inches(1.45), "16.7 kg", "Limestone-equivalent alkalinity for 10 kg of pyrite.")

    add_panel(slide, Inches(0.8), Inches(3.65), Inches(5.7), Inches(2.8), "Set-up")
    add_bullets(
        slide,
        Inches(1.05),
        Inches(4.15),
        Inches(5.15),
        Inches(2.0),
        [
            "1 mol pyrite -> about 4 mol H+",
            "1 mol CaCO3 neutralizes 2 mol H+",
            "So 1 mol pyrite needs about 2 mol CaCO3",
            "Mass ratio = (2 x 100) / 120 = 1.67",
        ],
        font_size=16,
    )

    add_panel(slide, Inches(6.8), Inches(3.65), Inches(5.45), Inches(2.8), "Engineering meaning")
    add_bullets(
        slide,
        Inches(7.05),
        Inches(4.15),
        Inches(4.95),
        Inches(2.0),
        [
            "Even small sulfide contents can create large long-term treatment obligations.",
            "Waste-rock design and source control matter because neutralization demand scales quickly.",
            "The point is not exact site chemistry, but the order of magnitude of the liability.",
        ],
        font_size=16,
    )

    add_source(slide, "Sources: OSMRE (2000); Zipper et al. (2023).")
    add_footer(slide, number)


def case_study_slide(prs, number):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header(slide, "Iron Mountain Mine, California", "Case study")

    add_metric(slide, Inches(0.82), Inches(1.8), Inches(2.6), Inches(1.45), "pH -3.6", "Reported by USGS in Richmond Mine waters - one of the most extreme AMD measurements recorded.")
    add_metric(slide, Inches(3.7), Inches(1.8), Inches(2.6), Inches(1.45), "200 g/L metals", "Peak dissolved metal concentrations reported in mine workings.")
    add_metric(slide, Inches(6.58), Inches(1.8), Inches(2.6), Inches(1.45), "760 g/L sulfate", "Illustrates how chemically concentrated AMD can become.")
    add_metric(slide, Inches(9.46), Inches(1.8), Inches(2.6), Inches(1.45), "400 million gal/yr", "Average AMD treated each year at the Minnesota Flats Treatment Plant.")

    add_panel(slide, Inches(0.82), Inches(3.6), Inches(5.8), Inches(2.9), "Why this case matters")
    add_bullets(
        slide,
        Inches(1.07),
        Inches(4.05),
        Inches(5.3),
        Inches(2.1),
        [
            "EPA once described the site as releasing more than one ton of toxic metals per day.",
            "The Sacramento River watershed suffered major ecological damage before large-scale remediation.",
            "Iron Mountain shows how AMD becomes permanent infrastructure, not a short-term fix.",
        ],
        font_size=16,
    )

    add_panel(slide, Inches(6.9), Inches(3.6), Inches(5.6), Inches(2.9), "Management lesson")
    add_bullets(
        slide,
        Inches(7.15),
        Inches(4.05),
        Inches(5.1),
        Inches(2.1),
        [
            "The permanent treatment plant began operating in 1994 and later shifted to high-density sludge processing.",
            "EPA notes cleanup has involved treatment of more than 8 billion gallons of acidic drainage over time.",
            "Once AMD reaches this scale, cleanup requires long-term funding, operations, and oversight.",
        ],
        font_size=16,
    )

    add_source(slide, "Sources: Nordstrom and Alpers (1999); Nordstrom et al. (2000); U.S. EPA (2006) and Iron Mountain site summaries.")
    add_footer(slide, number)


def solutions_slide(prs, number):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header(slide, "Technological solutions and alternatives", "Solutions")

    panels = [
        (
            Inches(0.7),
            Inches(1.82),
            Inches(3.85),
            Inches(2.2),
            "Prevention at the source",
            [
                "Keep oxygen and water away from reactive sulfide materials.",
                "Use covers, diversion ditches, sealing, and better waste storage design.",
                "Best option when feasible because it stops AMD before it forms.",
            ],
        ),
        (
            Inches(4.74),
            Inches(1.82),
            Inches(3.85),
            Inches(2.2),
            "Active treatment",
            [
                "Add lime or limestone to raise pH and precipitate metals.",
                "Works for severe sites such as Iron Mountain.",
                "Reliable, but energy-, reagent-, and sludge-intensive.",
            ],
        ),
        (
            Inches(8.78),
            Inches(1.82),
            Inches(3.85),
            Inches(2.2),
            "Passive and biological systems",
            [
                "Examples: anoxic limestone drains, open channels, wetlands, vertical flow ponds.",
                "Sulfate-reducing bacteria can remove sulfate and metals while generating alkalinity.",
                "Best for suitable chemistry, flow, climate, and land area.",
            ],
        ),
    ]

    for left, top, width, height, title, bullets in panels:
        add_panel(slide, left, top, width, height, title)
        add_bullets(slide, left + Inches(0.2), top + Inches(0.48), width - Inches(0.38), height - Inches(0.6), bullets, font_size=14)

    add_panel(slide, Inches(1.4), Inches(4.35), Inches(10.5), Inches(1.9), "Bottom line")
    add_text_block(
        slide,
        Inches(1.72),
        Inches(4.85),
        Inches(9.9),
        Inches(0.9),
        "No single technology works everywhere. The strongest engineering strategy combines prediction, source control, and the right treatment system for site chemistry and flow.",
        font_size=18,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    add_source(slide, "Sources: Johnson and Hallberg (2005); Skousen et al. (2017); Rambabu et al. (2020); U.S. EPA AMD overview.")
    add_footer(slide, number)


def economics_slide(prs, number):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header(slide, "Economics: why AMD becomes a long-tail liability", "Economics")

    add_metric(slide, Inches(0.82), Inches(1.84), Inches(3.35), Inches(1.45), "1,543 km", "Stream length protected by the Pennsylvania treatment systems evaluated by Black and Weber.")
    add_metric(slide, Inches(4.54), Inches(1.84), Inches(3.35), Inches(1.45), "$5,720 / km / year", "Present-value lifetime treatment cost implied by that analysis.")
    add_metric(slide, Inches(8.26), Inches(1.84), Inches(3.35), Inches(1.45), "AMDTreat", "OSMRE cost-estimation tool used for mine-drainage abatement planning.")

    add_panel(slide, Inches(0.82), Inches(3.7), Inches(5.75), Inches(2.75), "Economic problem")
    add_bullets(
        slide,
        Inches(1.06),
        Inches(4.15),
        Inches(5.25),
        Inches(1.95),
        [
            "Treatment often continues long after revenue from the mine has ended.",
            "Costs include reagents, sludge disposal, monitoring, inspections, and plant operations.",
            "If not internalized early, cleanup costs shift from private operators to the public.",
        ],
        font_size=16,
    )

    add_panel(slide, Inches(6.85), Inches(3.7), Inches(5.45), Inches(2.75), "Economic lesson")
    add_bullets(
        slide,
        Inches(7.1),
        Inches(4.15),
        Inches(4.95),
        Inches(1.95),
        [
            "Treatment can still be cost-effective because stream restoration benefits are large.",
            "But prevention and early reclamation are usually more rational than decades of cleanup.",
            "AMD is therefore an environmental issue and a financial-assurance issue.",
        ],
        font_size=16,
    )

    add_source(slide, "Sources: Black and Weber (2024); Johnson and Hallberg (2005); OSMRE AMDTreat resources; U.S. EPA AMD overview.")
    add_footer(slide, number)


def regulation_slide(prs, number):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header(slide, "Laws and regulation", "Policy")

    add_panel(slide, Inches(0.75), Inches(1.82), Inches(5.9), Inches(4.9), "Core legal framework")
    add_bullets(
        slide,
        Inches(1.0),
        Inches(2.3),
        Inches(5.4),
        Inches(3.95),
        [
            "Under the Clean Water Act, point-source discharges to waters of the United States generally need NPDES permits.",
            "EPA says permits are issued by EPA or authorized states and usually run no more than five years before renewal.",
            "SMCRA is the main federal law for environmental effects of coal mining, with states often acting as the primary regulator after primacy approval.",
            "Legacy hardrock sites such as Iron Mountain may also involve CERCLA / Superfund and state water boards.",
        ],
        font_size=15,
    )

    add_panel(slide, Inches(6.9), Inches(1.82), Inches(5.4), Inches(2.15), "40 CFR Part 434, Subpart C")
    add_bullets(
        slide,
        Inches(7.15),
        Inches(2.28),
        Inches(4.9),
        Inches(1.35),
        [
            "Existing-source limits include Fe 7.0 mg/L daily max and 3.5 mg/L average.",
            "Mn 4.0 / 2.0 mg/L, TSS 70 / 35 mg/L, and pH between 6.0 and 9.0.",
        ],
        font_size=14,
    )

    add_panel(slide, Inches(6.9), Inches(4.15), Inches(5.4), Inches(2.57), "Iron Mountain enforcement history")
    add_bullets(
        slide,
        Inches(7.15),
        Inches(4.62),
        Inches(4.9),
        Inches(1.7),
        [
            "California issued waste-discharge requirements in 1977 and an NPDES permit in 1978.",
            "A Cease and Desist Order followed in 1979 after permit violations.",
            "This makes the regulation story concrete, not just theoretical oversight.",
        ],
        font_size=14,
    )

    add_source(slide, "Sources: U.S. EPA NPDES overview (2025); eCFR 40 CFR Part 434 (2026); OSMRE SMCRA overview; California State Water Resources Control Board (1982).")
    add_footer(slide, number)


def conclusion_slide(prs, number):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header(slide, "Conclusion", "Takeaway")

    add_panel(slide, Inches(0.82), Inches(1.86), Inches(11.7), Inches(2.25), "Main takeaway")
    add_text_block(
        slide,
        Inches(1.12),
        Inches(2.38),
        Inches(11.1),
        Inches(1.25),
        "Acid mine drainage shows that environmental harm in energy and industrial systems begins upstream, when mining exposes sulfide-bearing rock and creates a geochemical system that can persist for decades.",
        font_size=22,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    add_panel(slide, Inches(0.82), Inches(4.45), Inches(5.6), Inches(2.0), "What to emphasize when presenting")
    add_bullets(
        slide,
        Inches(1.08),
        Inches(4.92),
        Inches(5.1),
        Inches(1.3),
        [
            "AMD is both a chemistry problem and a long-term engineering management problem.",
            "Iron Mountain proves the cost of waiting until contamination is already severe.",
        ],
        font_size=15,
    )

    add_panel(slide, Inches(6.78), Inches(4.45), Inches(5.74), Inches(2.0), "Closing statement")
    add_bullets(
        slide,
        Inches(7.05),
        Inches(4.92),
        Inches(5.2),
        Inches(1.3),
        [
            "Sustainable mining has to combine chemistry, hydrology, design, economics, and regulation from the beginning.",
            "The goal is prevention first, then treatment and accountability where needed.",
        ],
        font_size=15,
    )

    add_source(slide, "Sources: Nordstrom and Alpers (1999); Johnson and Hallberg (2005); U.S. EPA AMD overview.")
    add_footer(slide, number)


def references_slide(prs, number):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header(slide, "Selected references", "Bibliography")

    references = [
        "Black, Katie Jo, and Jeremy G. Weber. Communications Earth & Environment, 2024.",
        "California State Water Resources Control Board. Iron Mountain Mines review petition, 1982.",
        "Johnson, D. Barrie, and Kevin B. Hallberg. Science of the Total Environment, 2005.",
        "Nordstrom, D. Kirk, and Charles N. Alpers. PNAS, 1999.",
        "Rambabu, K. et al. Environmental Science and Ecotechnology, 2020.",
        "Skousen, Jeff et al. Mine Water and the Environment, 2017.",
        "OSMRE AMD guidance and SMCRA overview.",
        "U.S. EPA AMD, NPDES, and Iron Mountain resources.",
        "U.S. Geological Survey AMD overview.",
        "Zipper, Carl E. et al. Virginia Tech AMD treatment guide, 2023.",
    ]

    add_panel(slide, Inches(0.75), Inches(1.82), Inches(12.0), Inches(4.95))
    add_bullets(slide, Inches(1.0), Inches(2.15), Inches(11.5), Inches(4.25), references, font_size=15)
    add_source(slide, "Use the bibliography in your paper for full citation formatting on the submitted assignment.")
    add_footer(slide, number)


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    prs.core_properties.title = "Acid Mine Drainage (AMD) Presentation"
    prs.core_properties.subject = "CEE 304 Group 7 slide deck"
    prs.core_properties.author = "OpenAI Cursor Agent"
    prs.core_properties.keywords = "acid mine drainage, AMD, CEE 304, Iron Mountain, mining"

    title_slide(prs)
    bullet_slide(
        prs,
        2,
        "Challenge",
        "What AMD is and why it matters",
        [
            "AMD forms when mining exposes sulfide minerals, especially pyrite, to oxygen and water.",
            "That reaction generates sulfuric acid and dissolves metals into water, creating low-pH, sulfate-rich discharge.",
            "The harm is not limited to one day of extraction: waste rock, tailings, and abandoned tunnels can keep reacting for years or decades.",
            "So AMD is not just a local water-pollution problem - it is a supply-chain problem built into mining itself.",
        ],
        "Sources: Johnson and Hallberg (2005); U.S. EPA AMD overview; OSMRE AMD resources.",
    )
    bullet_slide(
        prs,
        3,
        "Supply chain",
        "How AMD fits into the energy supply chain",
        [
            "Mining sits near the beginning of the supply chain for coal, metals, wires, and industrial infrastructure.",
            "Blasting, excavation, crushing, ore processing, and waste storage all increase exposure of sulfide-bearing rock.",
            "Water then flows through waste piles, tunnels, and tailings, carrying acidity and dissolved metals into nearby watersheds.",
            "After closure, the challenge shifts from extraction to long-term containment, treatment, and monitoring.",
        ],
        "Sources: U.S. EPA AMD overview; OSMRE AMD resources; Johnson and Hallberg (2005).",
    )
    chemistry_slide(prs, 4)
    calculation_slide(prs, 5)
    case_study_slide(prs, 6)
    solutions_slide(prs, 7)
    economics_slide(prs, 8)
    regulation_slide(prs, 9)
    conclusion_slide(prs, 10)
    references_slide(prs, 11)
    prs.save(OUTPUT_PATH)


if __name__ == "__main__":
    build_presentation()
    print(f"Wrote {OUTPUT_PATH}")
